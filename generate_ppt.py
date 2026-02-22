"""
PG Management System - First Project Review Presentation Generator
Generates a professional .pptx file using python-pptx

Author: Priyanshi Kantariya (22BCE141)
Date: February 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# ============================================================================
# CONFIGURATION
# ============================================================================

STUDENT_INFO = {
    "name": "Priyanshi Kantariya",
    "roll_number": "22BCE141",
    "program": "B.Tech CSE",
    "internal_guide": "Prof. Monika G Shah",
    "external_guide": "Vivek Brahmabhatt",
    "company": "Urukrama Innovations Pvt. Ltd",
    "project_title": "PG Management System",
    "university": "Nirma University",
}

# Colors (RGB)
COLORS = {
    "primary": RGBColor(0, 51, 102),      # Dark Blue
    "secondary": RGBColor(0, 102, 153),   # Medium Blue
    "accent": RGBColor(255, 153, 0),      # Orange
    "text_dark": RGBColor(51, 51, 51),    # Dark Gray
    "text_light": RGBColor(255, 255, 255), # White
    "bg_light": RGBColor(240, 248, 255),  # Light Blue
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_title_shape(slide, title_text, left=0.5, top=0.3, width=9, height=0.8, font_size=32):
    """Add a styled title to a slide"""
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.LEFT
    return title_box


def add_content_box(slide, content_text, left=0.5, top=1.3, width=9, height=5, font_size=18):
    """Add content text box to a slide"""
    content_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    lines = content_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text_dark"]
        p.space_after = Pt(8)
    return content_box


def add_bullet_points(slide, title, bullets, top=1.3):
    """Add bullet points to a slide"""
    add_title_shape(slide, title)
    
    content = ""
    for bullet in bullets:
        content += f"• {bullet}\n"
    
    add_content_box(slide, content, top=top)


def add_table_slide(slide, title, headers, rows, top=1.5):
    """Add a table to a slide"""
    add_title_shape(slide, title)
    
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(top), Inches(9), Inches(0.4 * num_rows)).table
    
    # Set column widths
    col_width = 9 / num_cols
    for i in range(num_cols):
        table.columns[i].width = Inches(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS["text_light"]
            paragraph.font.size = Pt(14)
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = COLORS["text_dark"]


def add_shape_box(slide, text, left, top, width, height, fill_color, text_color=None, font_size=12):
    """Add a rounded rectangle shape with text"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS["primary"]
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color or COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a connector arrow between shapes"""
    # Using a line shape as arrow
    connector = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(start_left),
        Inches(start_top),
        Inches(0.3),
        Inches(0.4)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = COLORS["secondary"]
    connector.line.color.rgb = COLORS["secondary"]
    return connector


# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # University name at top
    uni_box = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(10), Inches(0.5))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = "NIRMA UNIVERSITY"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Institute name
    inst_box = slide.shapes.add_textbox(Inches(0), Inches(0.7), Inches(10), Inches(0.4))
    tf = inst_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Institute of Technology"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Project title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(1.8), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = STUDENT_INFO["project_title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(2.7), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "First Review Presentation"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Student details table
    details = [
        f"Presented By: {STUDENT_INFO['name']}",
        f"Roll Number: {STUDENT_INFO['roll_number']}",
        f"Program: {STUDENT_INFO['program']}",
        "",
        f"Internal Guide: {STUDENT_INFO['internal_guide']}",
        f"External Guide: {STUDENT_INFO['external_guide']}",
        "",
        f"Company: {STUDENT_INFO['company']}",
    ]
    
    details_box = slide.shapes.add_textbox(Inches(0), Inches(3.5), Inches(10), Inches(2.5))
    tf = details_box.text_frame
    for i, detail in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(10), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "February 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER


def create_toc_slide(prs):
    """Slide 2: Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Table of Contents")
    
    toc_items = [
        "1. Objective of Project",
        "2. Scope of Project",
        "3. Scope of First Review",
        "4. Functional Requirements",
        "5. Non-Functional Requirements",
        "6. Tools and Technologies",
        "7. Architectural Block Diagram",
        "8. Flowchart",
        "9. UML Diagrams",
        "10. Target for Next Review",
        "11. References",
    ]
    
    content = "\n".join(toc_items)
    add_content_box(slide, content, font_size=20)


def create_objective_slide(prs):
    """Slide 3: Objective of Project"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bullets = [
        "Develop a comprehensive web-based PG (Paying Guest) Management System",
        "Digitize and streamline property management operations for PG owners",
        "Provide tenants with a self-service portal for bills, complaints, and profile management",
        "Enable public users to browse properties and request visits online",
        "Implement role-based access control for Admin, Tenant, and Visitor roles",
        "Create a scalable, maintainable, and user-friendly application",
        "Reduce manual paperwork and improve operational efficiency",
    ]
    
    add_bullet_points(slide, "Objective of Project", bullets)


def create_scope_slide(prs):
    """Slide 4: Scope of Project"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bullets = [
        "Multi-role User System: Admin, Tenant, and Public Visitor portals",
        "Property Management: Add, edit, and display PG property listings",
        "Tenant Management: Registration, profile management, room assignments",
        "Billing System: Utility readings, bill generation, payment tracking",
        "Complaint Management: Raise, track, and resolve tenant complaints",
        "Visit Request System: Allow visitors to schedule property visits",
        "Authentication & Authorization: Secure login with role-based access",
        "Cloud-based Architecture: Firebase for backend services",
    ]
    
    add_bullet_points(slide, "Scope of Project", bullets)


def create_first_review_scope_slide(prs):
    """Slide 5: Scope of First Review"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "Scope of First Review")
    
    # Completed items
    completed = [
        "Project setup with React 19 + Vite",
        "Firebase integration (Auth, Firestore, Storage)",
        "User authentication (Admin & Tenant login/signup)",
        "Public portal with property listings",
        "Admin dashboard with property management",
        "Tenant dashboard with bill viewing",
        "Complaint management module",
        "Visit request functionality",
        "Responsive UI with TailwindCSS",
    ]
    
    content = "✓ Completed in Review 1:\n\n"
    for item in completed:
        content += f"  • {item}\n"
    
    add_content_box(slide, content, font_size=16)


def create_functional_requirements_slide(prs):
    """Slide 6: Functional Requirements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Functional Requirements")
    
    # Create two-column layout
    left_reqs = [
        "FR1: User Registration & Login",
        "FR2: Role-based Authentication",
        "FR3: View Property Listings",
        "FR4: Property Details View",
        "FR5: Request Property Visit",
        "FR6: Admin Dashboard Access",
        "FR7: Add/Edit Properties",
        "FR8: Manage Tenants",
        "FR9: Enter Utility Readings",
        "FR10: Generate Bills",
        "FR11: View Bills Overview",
        "FR12: Manage Visit Requests",
    ]
    
    right_reqs = [
        "FR13: Manage Complaints",
        "FR14: Reply to Complaints",
        "FR15: Resolve Complaints",
        "FR16: Tenant Dashboard Access",
        "FR17: View Personal Bills",
        "FR18: Report Payment",
        "FR19: View Own Complaints",
        "FR20: Create New Complaint",
        "FR21: View Tenant Profile",
        "FR22: Update Profile",
        "FR23: Secure Logout",
        "FR24: Session Management",
    ]
    
    # Left column
    left_content = "\n".join([f"• {r}" for r in left_reqs])
    add_content_box(slide, left_content, left=0.3, width=4.5, font_size=13)
    
    # Right column
    right_content = "\n".join([f"• {r}" for r in right_reqs])
    add_content_box(slide, right_content, left=5, width=4.5, font_size=13)


def create_nfr_slide(prs):
    """Slide 7: Non-Functional Requirements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "Non-Functional Requirements")
    
    headers = ["Category", "Requirement"]
    rows = [
        ["Performance", "Page load time < 3 seconds"],
        ["Scalability", "Support 100+ concurrent users"],
        ["Security", "Firebase Auth with secure rules"],
        ["Usability", "Mobile-responsive design"],
        ["Reliability", "99.9% uptime (Firebase SLA)"],
        ["Maintainability", "Modular component architecture"],
        ["Compatibility", "Chrome, Firefox, Safari, Edge"],
        ["Data Integrity", "Real-time sync with Firestore"],
    ]
    
    add_table_slide(slide, "", headers, rows, top=1.3)
    # Remove duplicate title
    slide.shapes[0]._element.getparent().remove(slide.shapes[0]._element)


def create_tools_tech_slide(prs):
    """Slide 8: Tools and Technologies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "Tools and Technologies")
    
    headers = ["Category", "Technology", "Purpose"]
    rows = [
        ["Frontend", "React 19", "UI Components & State"],
        ["Build Tool", "Vite 7", "Fast Development Server"],
        ["Styling", "TailwindCSS 4", "Utility-first CSS"],
        ["Routing", "React Router 7", "Client-side Navigation"],
        ["Backend", "Firebase", "BaaS Platform"],
        ["Database", "Firestore", "NoSQL Real-time DB"],
        ["Auth", "Firebase Auth", "User Authentication"],
        ["Storage", "Cloud Storage", "File/Image Storage"],
        ["Hosting", "Vercel", "Deployment Platform"],
        ["Version Control", "Git/GitHub", "Source Code Management"],
    ]
    
    add_table_slide(slide, "", headers, rows, top=1.3)
    slide.shapes[0]._element.getparent().remove(slide.shapes[0]._element)


def create_architecture_slide(prs):
    """Slide 9: Architecture Block Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Architectural Block Diagram")
    
    # Layer colors
    layer_colors = {
        "user": RGBColor(227, 242, 253),
        "presentation": RGBColor(255, 248, 225),
        "application": RGBColor(243, 229, 245),
        "service": RGBColor(224, 247, 250),
        "backend": RGBColor(255, 235, 238),
        "data": RGBColor(232, 245, 233),
    }
    
    # User Layer
    add_shape_box(slide, "USERS\n(Public | Tenant | Admin)", 3, 1.2, 4, 0.6, layer_colors["user"], COLORS["text_dark"], 11)
    
    # Arrow
    add_arrow(slide, 4.85, 1.85, 4.85, 2.1)
    
    # Presentation Layer
    add_shape_box(slide, "PRESENTATION LAYER", 0.5, 2.2, 9, 0.4, COLORS["secondary"], COLORS["text_light"], 12)
    add_shape_box(slide, "Public Portal\nHome | Properties | Visits", 0.5, 2.65, 2.8, 0.7, layer_colors["presentation"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Tenant Portal\nDashboard | Bills | Complaints", 3.5, 2.65, 2.8, 0.7, layer_colors["presentation"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Admin Portal\nProperties | Tenants | Billing", 6.5, 2.65, 2.8, 0.7, layer_colors["presentation"], COLORS["text_dark"], 10)
    
    # Arrow
    add_arrow(slide, 4.85, 3.4, 4.85, 3.6)
    
    # Application Layer
    add_shape_box(slide, "APPLICATION LAYER", 0.5, 3.7, 9, 0.4, COLORS["secondary"], COLORS["text_light"], 12)
    add_shape_box(slide, "React Router\nRoute Protection", 0.5, 4.15, 2.8, 0.6, layer_colors["application"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Auth Context\nState Management", 3.5, 4.15, 2.8, 0.6, layer_colors["application"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Shared Components\nUI Library", 6.5, 4.15, 2.8, 0.6, layer_colors["application"], COLORS["text_dark"], 10)
    
    # Arrow
    add_arrow(slide, 4.85, 4.8, 4.85, 5)
    
    # Service Layer
    add_shape_box(slide, "SERVICE LAYER (Firebase SDKs)", 0.5, 5.1, 9, 0.4, COLORS["secondary"], COLORS["text_light"], 12)
    
    # Arrow
    add_arrow(slide, 4.85, 5.55, 4.85, 5.75)
    
    # Backend Layer
    add_shape_box(slide, "FIREBASE BACKEND", 0.5, 5.85, 9, 0.4, COLORS["secondary"], COLORS["text_light"], 12)
    add_shape_box(slide, "Authentication", 0.5, 6.3, 2.8, 0.5, layer_colors["backend"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Firestore DB", 3.5, 6.3, 2.8, 0.5, layer_colors["backend"], COLORS["text_dark"], 10)
    add_shape_box(slide, "Cloud Storage", 6.5, 6.3, 2.8, 0.5, layer_colors["backend"], COLORS["text_dark"], 10)


def create_flowchart_slide(prs):
    """Slide 10: Flowchart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Application User Flow")
    
    # Simplified flowchart using shapes
    colors = {
        "start": RGBColor(200, 230, 201),
        "process": RGBColor(227, 242, 253),
        "decision": RGBColor(255, 243, 224),
        "admin": RGBColor(255, 235, 238),
        "tenant": RGBColor(232, 245, 233),
    }
    
    # Row 1: Start
    add_shape_box(slide, "User Visits\nWebsite", 4, 1.2, 2, 0.6, colors["start"], COLORS["text_dark"], 10)
    add_arrow(slide, 4.85, 1.85, 4.85, 2)
    
    # Row 2: Home
    add_shape_box(slide, "Home Page", 4, 2.1, 2, 0.5, colors["process"], COLORS["text_dark"], 10)
    add_arrow(slide, 4.85, 2.65, 4.85, 2.8)
    
    # Row 3: Decision
    decision = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(3.5), Inches(2.9), Inches(3), Inches(0.8))
    decision.fill.solid()
    decision.fill.fore_color.rgb = colors["decision"]
    decision.line.color.rgb = COLORS["secondary"]
    tf = decision.text_frame
    p = tf.paragraphs[0]
    p.text = "Select Action"
    p.font.size = Pt(10)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Branches
    add_shape_box(slide, "Browse\nProperties", 0.5, 4, 1.8, 0.7, colors["process"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Request\nVisit", 2.5, 4, 1.5, 0.7, colors["process"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Tenant\nLogin", 4.2, 4, 1.5, 0.7, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Admin\nLogin", 6, 4, 1.5, 0.7, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Signup", 7.7, 4, 1.5, 0.7, colors["process"], COLORS["text_dark"], 9)
    
    # Tenant Flow
    add_shape_box(slide, "Tenant Dashboard", 3, 5, 2, 0.5, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Bills | Complaints | Profile", 3, 5.6, 2, 0.4, colors["tenant"], COLORS["text_dark"], 8)
    
    # Admin Flow
    add_shape_box(slide, "Admin Dashboard", 5.5, 5, 2, 0.5, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Properties | Tenants | Billing", 5.5, 5.6, 2, 0.4, colors["admin"], COLORS["text_dark"], 8)
    
    # Legend
    legend_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(9), Inches(0.5))
    tf = legend_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Legend: 🟢 Start/End  🔵 Process  🟡 Decision  🟠 Admin Flow  🟢 Tenant Flow"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_dark"]


def create_uml_use_case_slide(prs):
    """Slide 11a: Use Case Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "UML: Use Case Diagram", font_size=28)
    
    colors = {
        "visitor": RGBColor(227, 242, 253),
        "tenant": RGBColor(232, 245, 233),
        "admin": RGBColor(255, 243, 224),
        "system": RGBColor(250, 250, 250),
    }
    
    # System boundary
    system = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.3), Inches(6), Inches(5.3))
    system.fill.solid()
    system.fill.fore_color.rgb = colors["system"]
    system.line.color.rgb = COLORS["primary"]
    
    # System title
    sys_title = slide.shapes.add_textbox(Inches(2), Inches(1.35), Inches(6), Inches(0.3))
    tf = sys_title.text_frame
    p = tf.paragraphs[0]
    p.text = "PG Management System"
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Actors (stick figures represented as ovals with labels)
    add_shape_box(slide, "👤\nVisitor", 0.3, 2, 1.3, 0.8, colors["visitor"], COLORS["text_dark"], 10)
    add_shape_box(slide, "👤\nTenant", 0.3, 3.5, 1.3, 0.8, colors["tenant"], COLORS["text_dark"], 10)
    add_shape_box(slide, "👤\nAdmin", 0.3, 5, 1.3, 0.8, colors["admin"], COLORS["text_dark"], 10)
    
    # Use Cases - organized in groups
    # Public
    add_shape_box(slide, "View Properties", 2.3, 1.8, 1.7, 0.4, RGBColor(225, 245, 254), COLORS["text_dark"], 9)
    add_shape_box(slide, "Request Visit", 4.2, 1.8, 1.5, 0.4, RGBColor(225, 245, 254), COLORS["text_dark"], 9)
    add_shape_box(slide, "Signup", 5.9, 1.8, 1.2, 0.4, RGBColor(225, 245, 254), COLORS["text_dark"], 9)
    
    # Auth
    add_shape_box(slide, "Login", 2.3, 2.4, 1.2, 0.4, RGBColor(243, 229, 245), COLORS["text_dark"], 9)
    add_shape_box(slide, "Logout", 3.7, 2.4, 1.2, 0.4, RGBColor(243, 229, 245), COLORS["text_dark"], 9)
    
    # Tenant Features
    add_shape_box(slide, "View Bills", 2.3, 3.2, 1.3, 0.4, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Report Payment", 3.8, 3.2, 1.5, 0.4, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "View Complaints", 5.5, 3.2, 1.6, 0.4, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Create Complaint", 2.3, 3.7, 1.6, 0.4, colors["tenant"], COLORS["text_dark"], 9)
    add_shape_box(slide, "View Profile", 4.1, 3.7, 1.4, 0.4, colors["tenant"], COLORS["text_dark"], 9)
    
    # Admin Features
    add_shape_box(slide, "Manage Properties", 2.3, 4.5, 1.8, 0.4, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Manage Tenants", 4.3, 4.5, 1.6, 0.4, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Enter Utilities", 6.1, 4.5, 1.5, 0.4, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Bills Overview", 2.3, 5, 1.5, 0.4, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Manage Visits", 4, 5, 1.5, 0.4, colors["admin"], COLORS["text_dark"], 9)
    add_shape_box(slide, "Manage Complaints", 5.7, 5, 1.8, 0.4, colors["admin"], COLORS["text_dark"], 9)
    
    # Legend
    legend = slide.shapes.add_textbox(Inches(2), Inches(6.65), Inches(6), Inches(0.3))
    tf = legend.text_frame
    p = tf.paragraphs[0]
    p.text = "→ Association  |  - - -> Include Relationship"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER


def create_uml_class_slide(prs):
    """Slide 11b: Class Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "UML: Class Diagram", font_size=28)
    
    def draw_class(slide, name, attributes, methods, left, top, width=2.2, color=RGBColor(227, 242, 253)):
        """Draw a UML class box"""
        height = 0.35 + (len(attributes) + len(methods)) * 0.18
        
        # Class name box
        name_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.35))
        name_shape.fill.solid()
        name_shape.fill.fore_color.rgb = COLORS["primary"]
        tf = name_shape.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Attributes box
        attr_height = len(attributes) * 0.18 + 0.05
        attr_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.35), Inches(width), Inches(attr_height))
        attr_shape.fill.solid()
        attr_shape.fill.fore_color.rgb = color
        tf = attr_shape.text_frame
        tf.word_wrap = True
        for i, attr in enumerate(attributes):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = attr
            p.font.size = Pt(8)
            p.font.color.rgb = COLORS["text_dark"]
        
        # Methods box
        meth_height = len(methods) * 0.18 + 0.05
        meth_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.35 + attr_height), Inches(width), Inches(meth_height))
        meth_shape.fill.solid()
        meth_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = meth_shape.text_frame
        tf.word_wrap = True
        for i, meth in enumerate(methods):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = meth
            p.font.size = Pt(8)
            p.font.color.rgb = COLORS["text_dark"]
    
    # Draw classes
    draw_class(slide, "User", 
               ["- id: string", "- email: string", "- role: string"],
               ["+ login()", "+ logout()"],
               0.3, 1.3)
    
    draw_class(slide, "Admin",
               ["- name: string", "- phone: string"],
               ["+ manageProperties()", "+ manageTenants()"],
               0.3, 3.2)
    
    draw_class(slide, "Tenant",
               ["- name: string", "- room: string", "- propertyId: string"],
               ["+ viewBills()", "+ createComplaint()"],
               2.8, 3.2)
    
    draw_class(slide, "Property",
               ["- id: string", "- name: string", "- address: string", "- rooms: number"],
               ["+ getDetails()", "+ getAvailability()"],
               5.3, 1.3)
    
    draw_class(slide, "Bill",
               ["- id: string", "- tenantId: string", "- amount: number", "- dueDate: date"],
               ["+ calculate()", "+ markPaid()"],
               5.3, 3.5)
    
    draw_class(slide, "Complaint",
               ["- id: string", "- tenantId: string", "- status: string"],
               ["+ submit()", "+ resolve()"],
               7.6, 1.3)
    
    draw_class(slide, "VisitRequest",
               ["- id: string", "- name: string", "- phone: string", "- date: date"],
               ["+ submit()", "+ updateStatus()"],
               7.6, 3.5)
    
    # Relationships notation
    rel_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.2), Inches(9), Inches(0.6))
    tf = rel_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Relationships: User ◁── Admin, Tenant (Inheritance) | Property ◇── Bill, Complaint (Composition)"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_dark"]


def create_uml_sequence_slide(prs):
    """Slide 11c: Sequence Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "UML: Sequence Diagram - Tenant Bill Payment", font_size=24)
    
    # Actors/Objects at top
    objects = [
        ("Tenant", 1),
        ("UI", 2.5),
        ("AuthContext", 4),
        ("Firestore", 5.5),
        ("Bill Service", 7),
    ]
    
    for name, left in objects:
        # Object box
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.3), Inches(1.3), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["primary"]
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Lifeline
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.6), Inches(1.7), Inches(0.05), Inches(4.5))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["secondary"]
        line.line.color.rgb = COLORS["secondary"]
    
    # Messages (simplified representation)
    messages = [
        ("1. Login", 0.5, 2, 1.5),
        ("2. Verify Credentials", 2, 2.3, 1.5),
        ("3. Auth Token", 1.5, 2.6, -1.5),
        ("4. Fetch Bills", 1, 3, 2),
        ("5. Query Bills", 2.5, 3.3, 2),
        ("6. Bill Data", 2.5, 3.6, -2),
        ("7. Display Bills", 1, 3.9, -1),
        ("8. Report Payment", 1, 4.3, 2),
        ("9. Update Payment", 2.5, 4.6, 2),
        ("10. Confirmation", 2.5, 4.9, -2),
        ("11. Show Success", 1, 5.2, -1),
    ]
    
    for msg, left, top, direction in messages:
        # Message label
        label = slide.shapes.add_textbox(Inches(left + 0.5), Inches(top - 0.15), Inches(2), Inches(0.25))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = msg
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS["text_dark"]
        
        # Arrow (simplified as line)
        arrow_width = abs(direction) * 0.8
        arrow_left = left + 1 if direction > 0 else left + 1 + direction * 0.8
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if direction > 0 else MSO_SHAPE.LEFT_ARROW,
                                       Inches(arrow_left), Inches(top + 0.05), Inches(arrow_width), Inches(0.1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS["accent"]


def create_next_review_slide(prs):
    """Slide 12: Target for Next Review"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bullets = [
        "Complete utility reading and auto-bill generation feature",
        "Implement payment gateway integration",
        "Add email/SMS notifications for bills and complaints",
        "Develop reporting and analytics dashboard for Admin",
        "Implement data export functionality (PDF bills, reports)",
        "Add advanced search and filtering capabilities",
        "Conduct unit and integration testing",
        "Performance optimization and security audit",
        "User Acceptance Testing (UAT) preparation",
        "Documentation completion",
    ]
    
    add_bullet_points(slide, "Target for Next Review (Review 2)", bullets)


def create_references_slide(prs):
    """Slide 13: References"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "References")
    
    references = [
        "[1] React Documentation. (2026). React – A JavaScript library for building user interfaces. https://react.dev/",
        "",
        "[2] Firebase Documentation. (2026). Firebase Documentation. https://firebase.google.com/docs",
        "",
        "[3] TailwindCSS. (2026). A utility-first CSS framework. https://tailwindcss.com/",
        "",
        "[4] Vite. (2026). Next Generation Frontend Tooling. https://vitejs.dev/",
        "",
        "[5] M. Fowler, \"Patterns of Enterprise Application Architecture,\" Addison-Wesley, 2002.",
        "",
        "[6] E. Gamma et al., \"Design Patterns: Elements of Reusable Object-Oriented Software,\" Addison-Wesley, 1994.",
        "",
        "[7] Nielsen, J. (1994). \"Usability Engineering.\" Morgan Kaufmann Publishers.",
    ]
    
    add_content_box(slide, "\n".join(references), font_size=12, top=1.2)


def create_thank_you_slide(prs):
    """Slide 14: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Thank You text
    thank_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(3.5), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0), Inches(5), Inches(10), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{STUDENT_INFO['name']} | {STUDENT_INFO['roll_number']}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["text_dark"]
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = STUDENT_INFO['company']
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER


# ============================================================================
# MAIN FUNCTION
# ============================================================================

def generate_presentation():
    """Generate the complete presentation"""
    print("🚀 Starting PPT Generation...")
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Generate all slides
    print("📝 Creating Title Slide...")
    create_title_slide(prs)
    
    print("📋 Creating Table of Contents...")
    create_toc_slide(prs)
    
    print("🎯 Creating Objective Slide...")
    create_objective_slide(prs)
    
    print("📐 Creating Scope Slide...")
    create_scope_slide(prs)
    
    print("✅ Creating First Review Scope...")
    create_first_review_scope_slide(prs)
    
    print("⚙️ Creating Functional Requirements...")
    create_functional_requirements_slide(prs)
    
    print("🔒 Creating Non-Functional Requirements...")
    create_nfr_slide(prs)
    
    print("🛠️ Creating Tools & Technologies...")
    create_tools_tech_slide(prs)
    
    print("🏗️ Creating Architecture Diagram...")
    create_architecture_slide(prs)
    
    print("🔄 Creating Flowchart...")
    create_flowchart_slide(prs)
    
    print("📊 Creating Use Case Diagram...")
    create_uml_use_case_slide(prs)
    
    print("📦 Creating Class Diagram...")
    create_uml_class_slide(prs)
    
    print("🔀 Creating Sequence Diagram...")
    create_uml_sequence_slide(prs)
    
    print("🎯 Creating Next Review Targets...")
    create_next_review_slide(prs)
    
    print("📚 Creating References...")
    create_references_slide(prs)
    
    print("🙏 Creating Thank You Slide...")
    create_thank_you_slide(prs)
    
    # Save the presentation
    output_file = "PG_Management_System_Review1.pptx"
    prs.save(output_file)
    
    print(f"\n✅ Presentation generated successfully!")
    print(f"📁 Output file: {os.path.abspath(output_file)}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return output_file


if __name__ == "__main__":
    generate_presentation()
